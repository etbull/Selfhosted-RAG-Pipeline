import os
import json
import requests 
from dotenv import load_dotenv 
from azure.storage.blob import BlobServiceClient
from datetime import datetime, timezone, timedelta
from typing import Optional, Dict, Any, List
from functools import lru_cache

# For text extraction
import pdfplumber
from docx import Document
from pptx import Presentation
from io import BytesIO 
import openpyxl
import xlrd


# Load .env
load_dotenv()

TENANT_ID = os.getenv("TENANT_ID")
CLIENT_ID = os.getenv("SHAREPOINT_CLIENT_ID")
CLIENT_SECRET = os.getenv("SHAREPOINT_CLIENT_SECRET")

# Option 1: Auto-discover sites (set to True)
# Option 2: Manual site list (set to False and populate SITES_LIST)
AUTO_DISCOVER_SITES = True

# Manual site list (only used if AUTO_DISCOVER_SITES = False)
# Format: site ID from Graph API
SITES_LIST = [
    # "yourdomain.sharepoint.com,site-guid,web-guid",
]

# Cache token for 55 minutes (tokens expire after 60 minutes)
_token_cache = {"token": None, "expires_at": None}

def get_graph_token():
    """Generate Graph App-only token with caching"""
    now = datetime.now(timezone.utc)
    
    # Return cached token if still valid
    if _token_cache["token"] and _token_cache["expires_at"] and now < _token_cache["expires_at"]:
        return _token_cache["token"]
    
    token_url = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"
    data = {
        "grant_type": "client_credentials",
        "client_id": CLIENT_ID,
        "client_secret": CLIENT_SECRET,
        "scope": "https://graph.microsoft.com/.default"
    }
    r = requests.post(token_url, data=data)
    r.raise_for_status()
    token_data = r.json()
    
    # Cache token (expires in 3600 seconds, cache for 3300 to be safe)
    _token_cache["token"] = token_data["access_token"]
    _token_cache["expires_at"] = now.replace(microsecond=0) + timedelta(seconds=3300)
    
    return _token_cache["token"]

def discover_all_sites() -> List[str]:
    """Auto-discover all SharePoint sites in the tenant"""
    headers = {"Authorization": f"Bearer {get_graph_token()}"}
    sites = []
    
    print("🔍 Discovering SharePoint sites...")
    
    # Get all sites
    url = "https://graph.microsoft.com/v1.0/sites?search=*"
    
    while url:
        r = requests.get(url, headers=headers)
        r.raise_for_status()
        data = r.json()
        
        for site in data.get("value", []):
            # Extract site path in the format needed
            site_id = site.get("id", "")
            web_url = site.get("webUrl", "")
            display_name = site.get("displayName", "")
            
            # Skip root site and personal sites
            if "/sites/" in web_url or "/teams/" in web_url:
                sites.append(site_id)
                print(f"  ✓ Found: {display_name} ({site_id})")
        
        # Handle pagination
        url = data.get("@odata.nextLink")
    
    print(f"📊 Total sites discovered: {len(sites)}\n")
    return sites

def get_site_info(site_id: str) -> Dict[str, Any]:
    """Fetch site information for a given site ID"""
    headers = {"Authorization": f"Bearer {get_graph_token()}"}
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}"
    r = requests.get(url, headers=headers)
    r.raise_for_status()
    site_data = r.json()
    
    return {
        "siteName": site_data.get("displayName", "Unknown"),
        "siteId": site_data.get("id", "Unknown"),
        "webUrl": site_data.get("webUrl", "Unknown")
    }

def get_all_document_libraries(site_id: str) -> List[Dict[str, Any]]:
    """Get all document libraries (drives) for a site"""
    headers = {"Authorization": f"Bearer {get_graph_token()}"}
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives"
    
    libraries = []
    
    try:
        r = requests.get(url, headers=headers)
        r.raise_for_status()
        drives = r.json().get("value", [])
        
        for drive in drives:
            libraries.append({
                "id": drive.get("id"),
                "name": drive.get("name"),
                "driveType": drive.get("driveType"),
                "webUrl": drive.get("webUrl")
            })
            
    except Exception as e:
        print(f"⚠️ Error fetching libraries for site {site_id}: {e}")
    
    return libraries

def get_file_permissions(site_id: str, drive_id: str, item_id: str) -> Dict[str, Any]:
    """Get permissions for a specific file/folder"""
    headers = {"Authorization": f"Bearer {get_graph_token()}"}
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{item_id}/permissions"
    
    try:
        r = requests.get(url, headers=headers)
        r.raise_for_status()
        permissions = r.json().get("value", [])
        
        allowed_users = []
        allowed_groups = []
        is_public = False
        
        for perm in permissions:
            # Check if it's a link (shared with anyone)
            if "link" in perm and perm.get("link", {}).get("scope") in ["anonymous", "organization"]:
                is_public = True
                continue
            
            # Get user principals (UPNs)
            if "grantedToV2" in perm:
                granted = perm["grantedToV2"]
                
                # Individual user
                if "user" in granted:
                    user_email = granted["user"].get("email")
                    if user_email:
                        allowed_users.append(user_email.lower())
                
                # Group
                if "group" in granted:
                    group_id = granted["group"].get("id")
                    group_email = granted["group"].get("email")
                    if group_id:
                        allowed_groups.append(group_id)
                    if group_email:
                        allowed_groups.append(group_email.lower())
            
            # Legacy grantedTo format
            elif "grantedTo" in perm:
                granted = perm["grantedTo"]
                if "user" in granted:
                    user_email = granted["user"].get("email")
                    if user_email:
                        allowed_users.append(user_email.lower())
        
        # Handle inherited permissions (most common case)
        has_inherited = len(permissions) == 0 or any(
            perm.get("inheritedFrom") for perm in permissions
        )
        
        return {
            "allowedUsers": list(set(allowed_users)),  # Remove duplicates
            "allowedGroups": list(set(allowed_groups)),
            "hasInheritedPermissions": has_inherited,
            "isPublicWithinOrg": is_public
        }
        
    except Exception as e:
        print(f"⚠️ Error fetching permissions for item {item_id}: {e}")
        return {
            "allowedUsers": [],
            "allowedGroups": [],
            "hasInheritedPermissions": True,
            "isPublicWithinOrg": False,
            "permissionError": str(e)
        }

def extract_text(name: str, stream: BytesIO) -> Optional[str]:
    """Extract text from PDF, DOCX, PPTX, TXT, CSV, JSON, HTML, and Excel."""
    suffix = name.lower().split(".")[-1]
    stream.seek(0)

    try:
        # ---------------------- PDF ----------------------
        if suffix == "pdf":
            text = ""
            with pdfplumber.open(stream) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
            return text.strip()

        # ---------------------- DOCX ----------------------
        if suffix == "docx":
            doc = Document(stream)
            return "\n".join([p.text for p in doc.paragraphs]).strip()

        # ---------------------- PPTX ----------------------
        if suffix == "pptx":
            prs = Presentation(stream)
            slide_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            slide_text.append(p.text)
            return "\n".join(slide_text).strip()

        # ---------------------- CSV ----------------------
        if suffix == "csv":
            text = stream.read().decode("utf-8", errors="ignore")
            return text.strip()

        # ---------------------- Excel (.xlsx, .xlsm) ----------------------
        if suffix in ["xlsx", "xlsm"]:
            wb = openpyxl.load_workbook(stream, read_only=True, data_only=True)
            output = []

            for sheet in wb.sheetnames:
                ws = wb[sheet]
                output.append(f"[{sheet}]")
                for row in ws.iter_rows(values_only=True):
                    row_text = ",".join("" if v is None else str(v) for v in row)
                    output.append(row_text)
                output.append("")

            return "\n".join(output).strip()

        # ---------------------- Excel (.xls – legacy) ----------------------
        if suffix == "xls":
            try:
                stream.seek(0)
                book = xlrd.open_workbook(file_contents=stream.read())
                output = []

                for i in range(book.nsheets):
                    sheet = book.sheet_by_index(i)
                    output.append(f"[{sheet.name}]")
                    for r in range(sheet.nrows):
                        row = sheet.row_values(r)
                        row_text = ",".join("" if v is None else str(v) for v in row)
                        output.append(row_text)
                    output.append("")

                return "\n".join(output).strip()
            except Exception as e:
                print(f"⚠️ XLS extraction failed: {e}")

        # ---------------------- Simple text formats ----------------------
        if suffix in ["txt", "md", "json", "xml", "html"]:
            return stream.read().decode("utf-8", errors="ignore").strip()

    except Exception as e:
        print(f"⚠️ Text extraction failed for {name}: {e}")

    return None

def build_metadata(item: Dict[str, Any], current_path: str, extracted_text: Optional[str], 
                   site_info: Dict[str, Any], library_name: str) -> Dict[str, Any]:
    """Build comprehensive metadata object"""
    metadata = {
        # File identification
        "fileName": item["name"],
        "fileExtension": item["name"].split(".")[-1].lower() if "." in item["name"] else "",
        "fileId": item["id"],
        "sharePointPath": current_path,
        "mimeType": item.get("file", {}).get("mimeType", "unknown"),
        
        # Site information
        "siteName": site_info["siteName"],
        "siteId": site_info["siteId"],
        "siteUrl": site_info["webUrl"],
        "libraryName": library_name,
        
        # Author/creator information
        "createdBy": item.get("createdBy", {}).get("user", {}).get("displayName", "Unknown"),
        "createdByEmail": item.get("createdBy", {}).get("user", {}).get("email", "Unknown"),
        "createdAt": item.get("createdDateTime", "Unknown"),
        
        # Last modifier information
        "lastModifiedBy": item.get("lastModifiedBy", {}).get("user", {}).get("displayName", "Unknown"),
        "lastModifiedByEmail": item.get("lastModifiedBy", {}).get("user", {}).get("email", "Unknown"),
        "lastModifiedAt": item["lastModifiedDateTime"],
        
        # File properties
        "size": item.get("size", 0),
        "sizeReadable": format_bytes(item.get("size", 0)),
        "webUrl": item.get("webUrl"),
        
        # Sync metadata
        "syncedAt": datetime.now(timezone.utc).isoformat(),
        "hasExtractedText": extracted_text is not None and len(extracted_text) > 0,
    }
    
    # Add file hash if available
    if "file" in item and "hashes" in item["file"]:
        hashes = item["file"]["hashes"]
        if "quickXorHash" in hashes:
            metadata["quickXorHash"] = hashes["quickXorHash"]
        if "sha1Hash" in hashes:
            metadata["sha1Hash"] = hashes["sha1Hash"]
    
    # Add extracted text if available
    if extracted_text:
        metadata["contentText"] = extracted_text
        metadata["contentTextLength"] = len(extracted_text)
    
    return metadata

def format_bytes(size: int) -> str:
    """Convert bytes to human-readable format"""
    for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
        if size < 1024.0:
            return f"{size:.2f} {unit}"
        size /= 1024.0
    return f"{size:.2f} PB"

def sanitize_folder_name(name: str) -> str:
    """Sanitize site/library names for blob storage paths"""
    # Remove or replace invalid characters for blob paths
    invalid_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*']
    for char in invalid_chars:
        name = name.replace(char, '_')
    return name.strip()

def sync_library_children(site_id: str, drive_id: str, item_id: Optional[str] = None, 
                          relative_path: str = "", site_info: Dict[str, Any] = None,
                          library_name: str = "Documents", blob_base_path: str = "",
                          library_permissions: Optional[Dict[str, Any]] = None):
    """Recursively sync all items in a document library"""
    
    # Build URL
    if item_id:
        url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{item_id}/children"
    else:
        url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/root/children"

    headers = {"Authorization": f"Bearer {get_graph_token()}"}
    
    try:
        r = requests.get(url, headers=headers)
        r.raise_for_status()
        items = r.json().get("value", [])
    except Exception as e:
        print(f"⚠️ Error accessing library: {e}")
        return

    for item in items:
        current_path = f"{relative_path}/{item['name']}".lstrip("/")

        # Folder → recurse
        if "folder" in item:
            print(f"📁 Folder: {blob_base_path}/{current_path}")
            sync_library_children(site_id, drive_id, item["id"], current_path, 
                                 site_info, library_name, blob_base_path, library_permissions)

        # File → check timestamps, upload if changed
        else:
            # Confirms that the directory to download is there
            directoryName = "RAG_DATA_ROOT"
            if not os.path.exists(directoryName):
                os.mkdir(directoryName)
                print(f"Directory made: {directoryName}")

            # Defining some required variables
            file_blob_path = f"{blob_base_path}/{current_path}"
            meta_blob_path = file_blob_path + ".meta.json"

            # Uploads data
            uploadData(file_blob_path, meta_blob_path, site_id, drive_id, item, headers, current_path, site_info, library_name, directoryName)

def uploadData(file_blob_path, meta_blob_path, site_id, drive_id, item, headers, current_path, site_info, library_name, directoryName):
    """Actually downloads data to the local folder for ingesting
    1. Checks if the file exists / has been changed
    2. If required, downloads the data
    """
    print(f"Uploading Data: {file_blob_path}")

    # Download file bytes from SharePoint
    file_resp = requests.get(
        f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{item['id']}/content",
        headers=headers
    )
    file_resp.raise_for_status()
    file_bytes = file_resp.content
    extracted_text = extract_text(item["name"], BytesIO(file_bytes))

    # Build comprehensive metadata
    metadata = build_metadata(item, current_path, extracted_text, 
                            site_info, library_name)

    # Upload metadata JSON
    name=meta_blob_path
    name=name[0].split('/')[-1]
    data=json.dumps(metadata, indent=2)

    with open(directoryName, 'w') as jsonFile:
        json.dump(data, jsonFile)
        print(f"Dumped: {name}")

def sync_site(site_id: str):
    """Sync all document libraries in a single site"""
    
    # Get site information
    site_info = get_site_info(site_id)
    site_name = sanitize_folder_name(site_info["siteName"])
    
    print(f"\n{'='*80}")
    print(f"Syncing site: {site_info['siteName']}")
    print(f"{'='*80}")
    
    # Get all document libraries for this site
    libraries = get_all_document_libraries(site_id)
    
    if not libraries:
        print(f"No document libraries found for {site_info['siteName']}")
        return
    
    print(f"Found {len(libraries)} document libraries:")
    for lib in libraries:
        print(f"  • {lib['name']} ({lib['driveType']})")
    
    # Sync each library
    for library in libraries:
        library_name = sanitize_folder_name(library["name"])
        blob_base_path = f"{site_name}/{library_name}"
        
        print(f"\nSyncing library: {library['name']}")
        print(f"   Blob path: {blob_base_path}/")
                
        try:
            sync_library_children(
                site_id=site_id,
                drive_id=library["id"],
                item_id=None,
                relative_path="",
                site_info=site_info,
                library_name=library["name"],
                blob_base_path=blob_base_path,
            )
        except Exception as e:
            print(f"❌ Error syncing library {library['name']}: {e}")
            continue

def sync_all_sites():
    """Main sync function - syncs all sites"""
    
    # Determine which sites to sync
    if AUTO_DISCOVER_SITES:
        sites = discover_all_sites()
    else:
        sites = SITES_LIST
        print(f"Using manual site list ({len(sites)} sites)")
    
    if not sites:
        print("No sites to sync!")
        return
    
    print(f"\nStarting sync for {len(sites)} site(s)...\n")
    
    # Sync each site
    for site_id in sites:
        try:
            sync_site(site_id)
        except Exception as e:
            print(f"Error syncing site {site_id}: {e}")
            continue
    
    print(f"\n{'='*80}")
    print("All sites sync completed")
    print(f"{'='*80}")

def get_user_groups_from_graph(user_upn: str) -> List[str]:
    """
    Get all groups a user belongs to (for expanding group permissions).
    This uses app-only token, so it can query any user's groups.
    """
    headers = {"Authorization": f"Bearer {get_graph_token()}"}
    url = f"https://graph.microsoft.com/v1.0/users/{user_upn}/transitiveMemberOf/microsoft.graph.group"
    
    groups = []
    try:
        r = requests.get(url, headers=headers)
        r.raise_for_status()
        
        for group in r.json().get("value", []):
            groups.append(group.get("id"))
            if group.get("mail"):
                groups.append(group.get("mail").lower())
        
        print(f"✓ User {user_upn} belongs to {len(groups)} groups")
        
    except Exception as e:
        print(f"⚠️ Error fetching groups for {user_upn}: {e}")
    
    return groups

# Start sync
if __name__ == "__main__":
    try:
        sync_all_sites()
    except Exception as e:
        print(f"❌ Sync failed: {e}")
        raise